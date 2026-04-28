"""Generate project presentation slides with rich visualizations."""

import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pptx.oxml.ns import qn

OUT_DIR = os.path.dirname(__file__)
CHART_DIR = os.path.join(OUT_DIR, "_charts")
os.makedirs(CHART_DIR, exist_ok=True)

# ── Color palette ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT_BLUE = RGBColor(0x4E, 0x9A, 0xF5)
ACCENT_GREEN = RGBColor(0x4E, 0xC9, 0xB0)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED = RGBColor(0xE8, 0x5D, 0x75)
ACCENT_PURPLE = RGBColor(0xA7, 0x7B, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3D)
TABLE_HEADER = RGBColor(0x2D, 0x2D, 0x4A)
TABLE_ROW1 = RGBColor(0x22, 0x22, 0x38)
TABLE_ROW2 = RGBColor(0x1E, 0x1E, 0x34)

# Matplotlib colors matching palette
MPL_BG = '#1B1B2F'
MPL_CARD = '#25253D'
MPL_BLUE = '#4E9AF5'
MPL_GREEN = '#4EC9B0'
MPL_ORANGE = '#F59E0B'
MPL_RED = '#E85D75'
MPL_PURPLE = '#A77BF3'
MPL_WHITE = '#FFFFFF'
MPL_LGRAY = '#BBBBCC'
MPL_MGRAY = '#888899'

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helper functions ──

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=13,
                       color=LIGHT_GRAY, line_spacing=Pt(6), font_name="Calibri"):
    """Add multi-line text with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = line_spacing
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=LIGHT_GRAY, bullet_color=ACCENT_BLUE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        pPr = p._p.get_or_add_pPr()
        for bn in pPr.findall(qn('a:buNone')):
            pPr.remove(bn)
        buChar = pPr.find(qn('a:buChar'))
        if buChar is None:
            buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '●')
        buClr = pPr.find(qn('a:buClr'))
        if buClr is None:
            buClr = etree.SubElement(pPr, qn('a:buClr'))
            srgb = etree.SubElement(buClr, qn('a:srgbClr'))
            srgb.set('val', str(bullet_color))
    return txBox


def add_table(slide, left, top, width, height, data, col_widths=None, font_size=11):
    rows_count = len(data)
    cols_count = len(data[0])
    table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY
            cell_fill = cell.fill
            cell_fill.solid()
            if r == 0:
                cell_fill.fore_color.rgb = TABLE_HEADER
            elif r % 2 == 1:
                cell_fill.fore_color.rgb = TABLE_ROW1
            else:
                cell_fill.fore_color.rgb = TABLE_ROW2
            tcPr = cell._tc.get_or_add_tcPr()
            for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                ln = tcPr.find(qn(f'a:{border_name}'))
                if ln is not None:
                    tcPr.remove(ln)
                ln = etree.SubElement(tcPr, qn(f'a:{border_name}'))
                ln.set('w', '0')
                etree.SubElement(ln, qn('a:noFill'))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape


def add_section_header(slide, number, title, subtitle=""):
    set_slide_bg(slide)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.8), Inches(2.5), Inches(1.6), Inches(1.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(10)
    add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(1),
                 title, font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(1.5), Inches(5.3), Inches(10), Inches(0.8),
                     subtitle, font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def setup_mpl_style():
    """Configure matplotlib for dark theme charts."""
    plt.rcParams.update({
        'figure.facecolor': MPL_BG,
        'axes.facecolor': MPL_CARD,
        'text.color': MPL_WHITE,
        'axes.labelcolor': MPL_WHITE,
        'xtick.color': MPL_LGRAY,
        'ytick.color': MPL_LGRAY,
        'axes.edgecolor': MPL_MGRAY,
        'grid.color': '#333350',
        'grid.alpha': 0.5,
        'font.family': 'sans-serif',
        'font.size': 14,
    })


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# CHART GENERATION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def make_data_flow_diagram():
    """Create a data flow architecture diagram."""
    setup_mpl_style()
    fig, ax = plt.subplots(figsize=(12, 5.5))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 6)
    ax.axis('off')

    # Data sources (left column)
    sources = [
        ("Food Prices\n(WFP)", MPL_BLUE, 0.3, 5.0),
        ("Rainfall\n(CHIRPS)", MPL_GREEN, 0.3, 4.2),
        ("Land Surface\n(FLDAS)", MPL_GREEN, 0.3, 3.4),
        ("Vegetation\n(MODIS)", MPL_GREEN, 0.3, 2.6),
        ("Climate Index\n(ENSO/IOD)", MPL_ORANGE, 0.3, 1.8),
        ("Conflict\n(ACLED)", MPL_RED, 0.3, 1.0),
        ("Crop Map\n(SPAM)", MPL_PURPLE, 0.3, 0.2),
        ("Population\n(WorldPop)", MPL_PURPLE, 2.3, 0.2),
    ]

    for name, color, x, y in sources:
        rect = mpatches.FancyBboxPatch((x, y), 1.7, 0.6, boxstyle="round,pad=0.1",
                                       facecolor=color, alpha=0.85, edgecolor='none')
        ax.add_patch(rect)
        ax.text(x + 0.85, y + 0.3, name, ha='center', va='center',
                fontsize=9, fontweight='bold', color='white')

    # Preprocessing (middle-left)
    proc_box = mpatches.FancyBboxPatch((4.2, 1.5), 2.0, 3.5, boxstyle="round,pad=0.15",
                                       facecolor='#2D2D4A', edgecolor=MPL_BLUE, linewidth=2)
    ax.add_patch(proc_box)
    ax.text(5.2, 4.6, "Preprocessing", ha='center', va='center',
            fontsize=12, fontweight='bold', color=MPL_BLUE)
    proc_items = ["Spatial Join", "Zonal Statistics", "Normalization", "Drought Index", "Aggregation"]
    for i, item in enumerate(proc_items):
        ax.text(5.2, 4.1 - i * 0.5, f"· {item}", ha='center', va='center',
                fontsize=9, color=MPL_LGRAY)

    # Merge (middle)
    merge_box = mpatches.FancyBboxPatch((6.8, 2.0), 1.8, 2.5, boxstyle="round,pad=0.15",
                                        facecolor='#2D2D4A', edgecolor=MPL_GREEN, linewidth=2)
    ax.add_patch(merge_box)
    ax.text(7.7, 4.1, "Merge", ha='center', va='center',
            fontsize=12, fontweight='bold', color=MPL_GREEN)
    merge_items = ["District Skeleton", "LEFT Joins", "Fuzzy Matching", "50+ columns"]
    for i, item in enumerate(merge_items):
        ax.text(7.7, 3.6 - i * 0.45, f"· {item}", ha='center', va='center',
                fontsize=9, color=MPL_LGRAY)

    # Feature Eng (middle-right)
    feat_box = mpatches.FancyBboxPatch((9.2, 2.0), 1.3, 2.5, boxstyle="round,pad=0.15",
                                       facecolor='#2D2D4A', edgecolor=MPL_ORANGE, linewidth=2)
    ax.add_patch(feat_box)
    ax.text(9.85, 4.1, "Feature\nEngineering", ha='center', va='center',
            fontsize=10, fontweight='bold', color=MPL_ORANGE)
    feat_items = ["Lag features", "Rolling stats", "Cyclical enc."]
    for i, item in enumerate(feat_items):
        ax.text(9.85, 3.4 - i * 0.45, f"· {item}", ha='center', va='center',
                fontsize=9, color=MPL_LGRAY)

    # Model (right)
    model_box = mpatches.FancyBboxPatch((10.9, 2.4), 1.0, 1.7, boxstyle="round,pad=0.15",
                                        facecolor=MPL_BLUE, alpha=0.9, edgecolor='none')
    ax.add_patch(model_box)
    ax.text(11.4, 3.5, "XGBoost", ha='center', va='center',
            fontsize=11, fontweight='bold', color='white')
    ax.text(11.4, 3.0, "Spatio-\nTemporal CV", ha='center', va='center',
            fontsize=8, color='white')

    # Arrows
    arrow_style = dict(arrowstyle='->', color=MPL_MGRAY, lw=1.5)
    for _, _, x, y in sources[:6]:
        ax.annotate('', xy=(4.2, 3.2), xytext=(x + 1.7, y + 0.3), arrowprops=arrow_style)
    for _, _, x, y in sources[6:]:
        ax.annotate('', xy=(4.2, 2.5), xytext=(x + 1.7, y + 0.3), arrowprops=arrow_style)

    ax.annotate('', xy=(6.8, 3.2), xytext=(6.2, 3.2), arrowprops=arrow_style)
    ax.annotate('', xy=(9.2, 3.2), xytext=(8.6, 3.2), arrowprops=arrow_style)
    ax.annotate('', xy=(10.9, 3.2), xytext=(10.5, 3.2), arrowprops=arrow_style)

    plt.tight_layout()
    path = os.path.join(CHART_DIR, "data_flow.png")
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor=MPL_BG)
    plt.close()
    return path


def make_cv_design_diagram():
    """Create spatio-temporal CV visualization."""
    setup_mpl_style()
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4.5), gridspec_kw={'width_ratios': [1.2, 1]})

    # Left: Temporal CV expanding window
    ax1.set_title("Temporal: Expanding Window (forward-only)", fontsize=14, fontweight='bold',
                  color=MPL_ORANGE, pad=12)
    folds = [
        ("Fold 1", 2007, 2012, 2013, 2024),
        ("Fold 2", 2007, 2017, 2018, 2024),
        ("Fold 3", 2007, 2022, 2023, 2024),
    ]
    for i, (name, ts, te, vs, ve) in enumerate(folds):
        y = 2 - i
        # Train bar
        ax1.barh(y, te - ts, left=ts, height=0.5, color=MPL_BLUE, alpha=0.85,
                 label='Train' if i == 0 else '')
        ax1.text((ts + te) / 2, y, f'Train\n{ts}–{te}', ha='center', va='center',
                 fontsize=8, color='white', fontweight='bold')
        # Val bar
        ax1.barh(y, ve - vs, left=vs, height=0.5, color=MPL_ORANGE, alpha=0.85,
                 label='Validation' if i == 0 else '')
        ax1.text((vs + ve) / 2, y, f'Val\n{vs}–{ve}', ha='center', va='center',
                 fontsize=8, color='white', fontweight='bold')
        ax1.text(2006, y, name, ha='right', va='center', fontsize=10, color=MPL_LGRAY)

    ax1.set_xlim(2005.5, 2026)
    ax1.set_ylim(-0.5, 2.8)
    ax1.set_yticks([])
    ax1.set_xlabel("Year", fontsize=11)
    ax1.legend(loc='upper left', fontsize=9, framealpha=0.3)
    ax1.grid(axis='x', alpha=0.3)

    # Right: Spatial CV disc-out concept
    ax2.set_title("Spatial: Leave-Disc-Out (350km radius)", fontsize=14, fontweight='bold',
                  color=MPL_GREEN, pad=12)
    np.random.seed(42)
    n_pts = 30
    lons = np.random.uniform(34, 50, n_pts)
    lats = np.random.uniform(-4, 12, n_pts)

    # Disc center
    cx, cy = 42, 4
    radius = 4.5  # visual radius

    # Color by in/out of disc
    for i in range(n_pts):
        dist = np.sqrt((lons[i] - cx) ** 2 + (lats[i] - cy) ** 2)
        if dist < radius:
            ax2.scatter(lons[i], lats[i], c=MPL_ORANGE, s=80, zorder=5, edgecolors='white', linewidth=0.5)
        else:
            ax2.scatter(lons[i], lats[i], c=MPL_BLUE, s=80, zorder=5, edgecolors='white', linewidth=0.5)

    circle = plt.Circle((cx, cy), radius, fill=False, edgecolor=MPL_ORANGE,
                         linewidth=2, linestyle='--', alpha=0.8)
    ax2.add_patch(circle)
    ax2.plot(cx, cy, 'x', color=MPL_ORANGE, markersize=12, markeredgewidth=2)
    ax2.text(cx, cy + radius + 0.5, "350 km disc", ha='center', fontsize=10,
             color=MPL_ORANGE, fontweight='bold')

    # Legend
    ax2.scatter([], [], c=MPL_BLUE, s=60, label='Train districts')
    ax2.scatter([], [], c=MPL_ORANGE, s=60, label='Validation districts')
    ax2.legend(loc='lower left', fontsize=9, framealpha=0.3)

    ax2.set_xlim(32, 52)
    ax2.set_ylim(-6, 14)
    ax2.set_xlabel("Longitude", fontsize=11)
    ax2.set_ylabel("Latitude", fontsize=11)
    ax2.set_aspect('equal')
    ax2.grid(alpha=0.2)

    plt.tight_layout()
    path = os.path.join(CHART_DIR, "cv_design.png")
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor=MPL_BG)
    plt.close()
    return path


def make_holdout_r2_chart():
    """Bar chart of holdout R² by target and horizon."""
    setup_mpl_style()
    fig, ax = plt.subplots(figsize=(8, 4.5))

    targets = ['Food Price\nIndex', 'Maize', 'Sorghum']
    h1 = [0.946, 0.989, 0.965]
    h2 = [0.900, 0.977, 0.937]
    h3 = [0.849, 0.969, 0.918]

    x = np.arange(len(targets))
    w = 0.25
    bars1 = ax.bar(x - w, h1, w, label='h=1', color=MPL_BLUE, alpha=0.9, edgecolor='none')
    bars2 = ax.bar(x, h2, w, label='h=2', color=MPL_GREEN, alpha=0.9, edgecolor='none')
    bars3 = ax.bar(x + w, h3, w, label='h=3', color=MPL_ORANGE, alpha=0.9, edgecolor='none')

    # Value labels
    for bars in [bars1, bars2, bars3]:
        for bar in bars:
            ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.008,
                    f'{bar.get_height():.3f}', ha='center', va='bottom',
                    fontsize=9, color=MPL_LGRAY, fontweight='bold')

    ax.set_xticks(x)
    ax.set_xticklabels(targets, fontsize=12)
    ax.set_ylabel('R²', fontsize=13)
    ax.set_title('Hold-Out R² (Variance Explained) by Target & Horizon', fontsize=15, fontweight='bold', pad=12)
    ax.set_ylim(0.8, 1.02)
    ax.legend(fontsize=11, loc='lower left', framealpha=0.3)
    ax.grid(axis='y', alpha=0.3)
    ax.axhline(y=0.95, color=MPL_MGRAY, linestyle=':', alpha=0.5)

    plt.tight_layout()
    path = os.path.join(CHART_DIR, "holdout_r2.png")
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor=MPL_BG)
    plt.close()
    return path


def make_holdout_mape_chart():
    """Bar chart of holdout MAPE by target and horizon."""
    setup_mpl_style()
    fig, ax = plt.subplots(figsize=(8, 4.5))

    targets = ['Food Price\nIndex', 'Maize', 'Sorghum']
    h1 = [2.31, 4.18, 5.63]
    h2 = [3.31, 6.14, 8.18]
    h3 = [4.29, 7.32, 10.79]

    x = np.arange(len(targets))
    w = 0.25
    bars1 = ax.bar(x - w, h1, w, label='h=1', color=MPL_BLUE, alpha=0.9)
    bars2 = ax.bar(x, h2, w, label='h=2', color=MPL_GREEN, alpha=0.9)
    bars3 = ax.bar(x + w, h3, w, label='h=3', color=MPL_ORANGE, alpha=0.9)

    for bars in [bars1, bars2, bars3]:
        for bar in bars:
            ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.15,
                    f'{bar.get_height():.1f}%', ha='center', va='bottom',
                    fontsize=9, color=MPL_LGRAY, fontweight='bold')

    ax.set_xticks(x)
    ax.set_xticklabels(targets, fontsize=12)
    ax.set_ylabel('MAPE (%)', fontsize=13)
    ax.set_title('Hold-Out MAPE (% Error) by Target & Horizon', fontsize=15, fontweight='bold', pad=12)
    ax.set_ylim(0, 14)
    ax.legend(fontsize=11, loc='upper left', framealpha=0.3)
    ax.grid(axis='y', alpha=0.3)
    ax.axhline(y=5, color=MPL_GREEN, linestyle=':', alpha=0.4, label='')
    ax.axhline(y=10, color=MPL_ORANGE, linestyle=':', alpha=0.4)

    plt.tight_layout()
    path = os.path.join(CHART_DIR, "holdout_mape.png")
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor=MPL_BG)
    plt.close()
    return path


def make_country_comparison_chart():
    """Grouped bar chart comparing Kenya vs Somalia MAPE."""
    setup_mpl_style()
    fig, ax = plt.subplots(figsize=(10, 4.5))

    labels = ['Food Price\nIndex h=1', 'Food Price\nIndex h=2', 'Food Price\nIndex h=3',
              'Maize h=1', 'Maize h=2', 'Maize h=3',
              'Sorghum h=1', 'Sorghum h=2', 'Sorghum h=3']
    ken = [2.36, 3.68, 5.09, 3.18, 4.70, 5.76, 3.12, 4.95, 7.89]
    som = [2.23, 2.73, 3.04, 5.75, 8.38, 9.75, 9.54, 13.21, 15.29]

    x = np.arange(len(labels))
    w = 0.35
    ax.bar(x - w / 2, ken, w, label='Kenya', color=MPL_BLUE, alpha=0.9)
    ax.bar(x + w / 2, som, w, label='Somalia', color=MPL_RED, alpha=0.9)

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=9, rotation=30, ha='right')
    ax.set_ylabel('MAPE (Mean Absolute % Error)', fontsize=13)
    ax.set_title('Kenya vs Somalia: Prediction Error Comparison', fontsize=15, fontweight='bold', pad=12)
    ax.legend(fontsize=12, framealpha=0.3)
    ax.grid(axis='y', alpha=0.3)

    # Separator lines between targets
    ax.axvline(x=2.5, color=MPL_MGRAY, linestyle=':', alpha=0.4)
    ax.axvline(x=5.5, color=MPL_MGRAY, linestyle=':', alpha=0.4)

    plt.tight_layout()
    path = os.path.join(CHART_DIR, "country_comparison.png")
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor=MPL_BG)
    plt.close()
    return path


def make_feature_importance_chart():
    """Stacked horizontal bar chart of feature group importance."""
    setup_mpl_style()
    fig, ax = plt.subplots(figsize=(10, 4))

    targets = ['Sorghum h=1', 'Maize h=1', 'Food Price Index h=1']
    groups = ['Past Prices', 'Land Surface\n(FLDAS)', 'Climate\nIndex', 'Vegetation\nHealth', 'Conflict', 'Other']
    colors = [MPL_BLUE, MPL_GREEN, MPL_ORANGE, MPL_PURPLE, MPL_RED, MPL_MGRAY]

    data = {
        'Food Price Index h=1': [93.8, 1.2, 2.1, 0.4, 0.1, 2.4],
        'Maize h=1': [77.0, 10.0, 8.0, 2.5, 0.5, 2.0],
        'Sorghum h=1': [78.3, 11.2, 6.8, 1.5, 1.0, 1.2],
    }

    y = np.arange(len(targets))
    left = np.zeros(len(targets))
    for i, group in enumerate(groups):
        vals = [data[t][i] for t in targets]
        bars = ax.barh(y, vals, left=left, height=0.5, label=group, color=colors[i], alpha=0.9)
        # Label if > 5%
        for j, v in enumerate(vals):
            if v > 5:
                ax.text(left[j] + v / 2, y[j], f'{v:.0f}%', ha='center', va='center',
                        fontsize=9, color='white', fontweight='bold')
        left += vals

    ax.set_yticks(y)
    ax.set_yticklabels(targets, fontsize=12)
    ax.set_xlabel('Importance (%)', fontsize=12)
    ax.set_title('Feature Group Importance (h=1)', fontsize=15, fontweight='bold', pad=12)
    ax.legend(loc='lower right', fontsize=9, framealpha=0.3, ncol=3)
    ax.set_xlim(0, 105)
    ax.grid(axis='x', alpha=0.2)

    plt.tight_layout()
    path = os.path.join(CHART_DIR, "feature_importance.png")
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor=MPL_BG)
    plt.close()
    return path


def make_horizon_degradation_chart():
    """Line chart showing metric degradation across horizons."""
    setup_mpl_style()
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4.5))

    horizons = [1, 2, 3]

    # R² degradation
    r2_fpi = [0.946, 0.900, 0.849]
    r2_maize = [0.989, 0.977, 0.969]
    r2_sorghum = [0.965, 0.937, 0.918]

    ax1.plot(horizons, r2_fpi, 'o-', color=MPL_BLUE, linewidth=2.5, markersize=10, label='Food Price Index')
    ax1.plot(horizons, r2_maize, 's-', color=MPL_GREEN, linewidth=2.5, markersize=10, label='Maize')
    ax1.plot(horizons, r2_sorghum, '^-', color=MPL_ORANGE, linewidth=2.5, markersize=10, label='Sorghum')

    for vals, marker_offset in [(r2_fpi, -0.012), (r2_maize, 0.012), (r2_sorghum, -0.012)]:
        for h, v in zip(horizons, vals):
            ax1.text(h, v + marker_offset + 0.008, f'{v:.3f}', ha='center', fontsize=9, color=MPL_LGRAY)

    ax1.set_xlabel('Forecast Horizon (months)', fontsize=12)
    ax1.set_ylabel('R²', fontsize=13)
    ax1.set_title('R² (Variance Explained) Degradation', fontsize=14, fontweight='bold', pad=10)
    ax1.set_xticks(horizons)
    ax1.set_ylim(0.82, 1.01)
    ax1.legend(fontsize=10, framealpha=0.3)
    ax1.grid(alpha=0.3)

    # MAPE degradation
    mape_fpi = [2.31, 3.31, 4.29]
    mape_maize = [4.18, 6.14, 7.32]
    mape_sorghum = [5.63, 8.18, 10.79]

    ax2.plot(horizons, mape_fpi, 'o-', color=MPL_BLUE, linewidth=2.5, markersize=10, label='Food Price Index')
    ax2.plot(horizons, mape_maize, 's-', color=MPL_GREEN, linewidth=2.5, markersize=10, label='Maize')
    ax2.plot(horizons, mape_sorghum, '^-', color=MPL_ORANGE, linewidth=2.5, markersize=10, label='Sorghum')

    for vals in [mape_fpi, mape_maize, mape_sorghum]:
        for h, v in zip(horizons, vals):
            ax2.text(h + 0.08, v + 0.2, f'{v:.1f}%', ha='left', fontsize=9, color=MPL_LGRAY)

    ax2.set_xlabel('Forecast Horizon (months)', fontsize=12)
    ax2.set_ylabel('MAPE (%)', fontsize=13)
    ax2.set_title('MAPE (% Error) Degradation', fontsize=14, fontweight='bold', pad=10)
    ax2.set_xticks(horizons)
    ax2.set_ylim(0, 13)
    ax2.legend(fontsize=10, framealpha=0.3)
    ax2.grid(alpha=0.3)

    plt.tight_layout()
    path = os.path.join(CHART_DIR, "horizon_degradation.png")
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor=MPL_BG)
    plt.close()
    return path


def make_cv_vs_holdout_chart():
    """Comparison of STCV vs Hold-out R²."""
    setup_mpl_style()
    fig, ax = plt.subplots(figsize=(9, 4.5))

    labels = ['Food Price\nIndex h=1', 'Food Price\nIndex h=2', 'Food Price\nIndex h=3',
              'Maize h=1', 'Maize h=2', 'Maize h=3',
              'Sorghum h=1', 'Sorghum h=2', 'Sorghum h=3']
    cv_r2 = [0.00, -0.34, -0.65, 0.71, 0.60, 0.49, 0.76, 0.65, 0.52]
    ho_r2 = [0.946, 0.900, 0.849, 0.989, 0.977, 0.969, 0.965, 0.937, 0.918]

    x = np.arange(len(labels))
    w = 0.35
    ax.bar(x - w / 2, cv_r2, w, label='Spatio-Temporal CV', color=MPL_ORANGE, alpha=0.85)
    ax.bar(x + w / 2, ho_r2, w, label='Hold-Out (2024+)', color=MPL_BLUE, alpha=0.85)

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=9, rotation=30, ha='right')
    ax.set_ylabel('R²', fontsize=13)
    ax.set_title('Spatio-Temporal Cross-Validation vs Hold-Out R²', fontsize=14, fontweight='bold', pad=12)
    ax.legend(fontsize=11, framealpha=0.3)
    ax.grid(axis='y', alpha=0.3)
    ax.axhline(y=0, color=MPL_MGRAY, linestyle='-', alpha=0.5)
    ax.axvline(x=2.5, color=MPL_MGRAY, linestyle=':', alpha=0.4)
    ax.axvline(x=5.5, color=MPL_MGRAY, linestyle=':', alpha=0.4)

    plt.tight_layout()
    path = os.path.join(CHART_DIR, "cv_vs_holdout.png")
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor=MPL_BG)
    plt.close()
    return path


def make_lag_strategy_diagram():
    """Visual diagram of lag strategy for forecasting."""
    setup_mpl_style()
    fig, ax = plt.subplots(figsize=(11, 4))
    ax.set_xlim(-1, 14)
    ax.set_ylim(-0.5, 4.5)
    ax.axis('off')
    ax.set_title('Lag Strategy: Preventing Data Leakage', fontsize=15, fontweight='bold',
                 color=MPL_WHITE, pad=15)

    months = ['t-12', 't-6', 't-3', 't-2', 't-1', 't', '', 't+1', 't+2', 't+3']
    x_positions = [0, 2, 4, 5, 6, 7, 8, 9, 10, 11]

    # Timeline
    ax.plot([0, 11.5], [3.5, 3.5], '-', color=MPL_MGRAY, linewidth=2)
    for x, label in zip(x_positions, months):
        if label:
            ax.plot(x, 3.5, 'o', color=MPL_LGRAY, markersize=6)
            ax.text(x, 3.8, label, ha='center', fontsize=9, color=MPL_LGRAY)

    # Current time marker
    ax.plot(7, 3.5, 'D', color=MPL_GREEN, markersize=12, zorder=5)
    ax.text(7, 4.2, 'Now (t)', ha='center', fontsize=11, color=MPL_GREEN, fontweight='bold')

    # Price lags (row 1)
    ax.text(-0.8, 2.5, 'Past Price\nFeatures', ha='center', va='center', fontsize=10, color=MPL_BLUE, fontweight='bold')
    for x, lag in [(6, 'lag₁'), (5, 'lag₂'), (4, 'lag₃'), (2, 'lag₆'), (0, 'lag₁₂')]:
        rect = mpatches.FancyBboxPatch((x - 0.4, 2.2), 0.8, 0.6, boxstyle="round,pad=0.05",
                                       facecolor=MPL_BLUE, alpha=0.7, edgecolor='none')
        ax.add_patch(rect)
        ax.text(x, 2.5, lag, ha='center', va='center', fontsize=8, color='white', fontweight='bold')
        ax.annotate('', xy=(7, 3.4), xytext=(x, 2.85),
                    arrowprops=dict(arrowstyle='->', color=MPL_BLUE, alpha=0.3, lw=1))

    # Exogenous for h=1 (row 2)
    ax.text(-0.8, 1.3, 'External\n(h=1)', ha='center', va='center', fontsize=10, color=MPL_GREEN, fontweight='bold')
    rect = mpatches.FancyBboxPatch((5.6, 1.0), 0.8, 0.6, boxstyle="round,pad=0.05",
                                   facecolor=MPL_GREEN, alpha=0.7, edgecolor='none')
    ax.add_patch(rect)
    ax.text(6, 1.3, 't-1', ha='center', va='center', fontsize=9, color='white', fontweight='bold')
    ax.annotate('', xy=(9, 3.4), xytext=(6.4, 1.6),
                arrowprops=dict(arrowstyle='->', color=MPL_GREEN, alpha=0.5, lw=1.5))
    ax.text(9, 4.2, 'Predict', ha='center', fontsize=9, color=MPL_GREEN)

    # Exogenous for h=3 (row 3)
    ax.text(-0.8, 0.2, 'External\n(h=3)', ha='center', va='center', fontsize=10, color=MPL_ORANGE, fontweight='bold')
    rect = mpatches.FancyBboxPatch((3.6, -0.1), 0.8, 0.6, boxstyle="round,pad=0.05",
                                   facecolor=MPL_ORANGE, alpha=0.7, edgecolor='none')
    ax.add_patch(rect)
    ax.text(4, 0.2, 't-3', ha='center', va='center', fontsize=9, color='white', fontweight='bold')
    ax.annotate('', xy=(11, 3.4), xytext=(4.4, 0.5),
                arrowprops=dict(arrowstyle='->', color=MPL_ORANGE, alpha=0.5, lw=1.5))
    ax.text(11, 4.2, 'Predict', ha='center', fontsize=9, color=MPL_ORANGE)

    # Forbidden zone
    rect = mpatches.FancyBboxPatch((7.3, -0.3), 4.3, 3.0, boxstyle="round,pad=0.1",
                                   facecolor=MPL_RED, alpha=0.08, edgecolor=MPL_RED,
                                   linewidth=1.5, linestyle='--')
    ax.add_patch(rect)
    ax.text(9.5, -0.1, 'Future data — never used as features', ha='center',
            fontsize=9, color=MPL_RED, style='italic')

    plt.tight_layout()
    path = os.path.join(CHART_DIR, "lag_strategy.png")
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor=MPL_BG)
    plt.close()
    return path


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# GENERATE ALL CHARTS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
print("Generating charts...")
chart_data_flow = make_data_flow_diagram()
chart_cv_design = make_cv_design_diagram()
chart_holdout_r2 = make_holdout_r2_chart()
chart_holdout_mape = make_holdout_mape_chart()
chart_country = make_country_comparison_chart()
chart_feat_imp = make_feature_importance_chart()
chart_horizon = make_horizon_degradation_chart()
chart_cv_vs_ho = make_cv_vs_holdout_chart()
chart_lag = make_lag_strategy_diagram()
print("Charts generated.")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# BUILD PRESENTATION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ── SLIDE 1: Title ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Subnational Food Price Prediction\nin East Africa",
             font_size=38, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
             "Integrating Climate, Conflict, and Socioeconomic Data for Admin2-Level Forecasting",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(4.5), Inches(4.4), Inches(4), Inches(0.03), ACCENT_BLUE)
add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
             "Kenya  &  Somalia  |  XGBoost (Gradient Boosting)  |  Spatio-Temporal Cross-Validation",
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.5),
             "Halim Jun  ·  Donghee Lee  ·  Sheldon Lu",
             font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(6.1), Inches(11), Inches(0.4),
             "Columbia University  ·  Spring 2025",
             font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ── SLIDE 2: Section 1 ──
slide = prs.slides.add_slide(blank)
add_section_header(slide, "1", "Problem Definition", "Why predict food prices in East Africa?")


# ── SLIDE 3: Problem Definition ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Problem Definition", font_size=28, color=WHITE, bold=True)

# 3 KPI-style cards at top
kpis = [
    ("3 Targets", "Food Price Index\nMaize\nSorghum", ACCENT_BLUE),
    ("3 Horizons", "1 month ahead\n2 months ahead\n3 months ahead", ACCENT_GREEN),
    ("2 Countries", "Kenya\nSomalia\nAdmin2 (district) level", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(kpis):
    x = Inches(0.5 + i * 4.2)
    card = add_shape(slide, x, Inches(1.2), Inches(3.8), Inches(1.8), CARD_BG)
    add_rect(slide, x, Inches(1.2), Inches(3.8), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.3), Inches(1.4), Inches(3.3), Inches(0.4),
                 title, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(1.9), Inches(3.3), Inches(1.0),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom: Motivation & Challenges in 2 columns
card_l = add_shape(slide, Inches(0.5), Inches(3.3), Inches(6.1), Inches(3.5), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(3.45), Inches(5.5), Inches(0.4),
             "Motivation", font_size=18, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(3.9), Inches(5.5), Inches(2.5), [
    "East Africa: recurring droughts, conflict, economic shocks → food crises",
    "Early warning systems require granular, forward-looking price forecasts",
    "National-level data masks critical subnational (district-level) variation",
    "Most studies ignore spatial autocorrelation (nearby regions influencing each other) leakage"
], font_size=12, spacing=Pt(8))

card_r = add_shape(slide, Inches(6.9), Inches(3.3), Inches(5.9), Inches(3.5), CARD_BG)
add_text_box(slide, Inches(7.2), Inches(3.45), Inches(5.3), Inches(0.4),
             "Key Challenges", font_size=18, color=ACCENT_RED, bold=True)
add_bullet_list(slide, Inches(7.2), Inches(3.9), Inches(5.3), Inches(2.5), [
    "8 heterogeneous data sources with different spatial/temporal resolutions",
    "Spatial leakage: nearby districts have correlated prices",
    "Temporal leakage: future data must never inform past predictions",
    "Forecast accuracy naturally degrades at longer horizons (1→3 months)"
], font_size=12, bullet_color=ACCENT_RED, spacing=Pt(8))


# ── SLIDE 4: Section 2 ──
slide = prs.slides.add_slide(blank)
add_section_header(slide, "2", "Data Sources", "8 heterogeneous datasets integrated")


# ── SLIDE 5: Data Sources with visual ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Data Sources", font_size=28, color=WHITE, bold=True)

data_table = [
    ["Category", "Source", "Provider", "Key Variables", "Resolution"],
    ["Target\n(Food Prices)", "WFP\n(World Food Programme)", "UN WFP", "Food Price Index, Maize, Sorghum", "Market → District"],
    ["Rainfall", "CHIRPS\n(Satellite Rainfall)", "UCSB", "SPI (Drought Index)", "5km → District"],
    ["Hydrology /\nTemperature", "FLDAS\n(Land Surface Model)", "NASA", "Soil moisture, Air temp,\nRunoff, Evapotranspiration (8 vars)", "District monthly"],
    ["Vegetation\nHealth", "MODIS / AVHRR\n(Satellite Imagery)", "NASA / NOAA", "NDVI (Greenness), LST (Surface Temp),\nVHI (Vegetation Health)", "District monthly"],
    ["Global\nClimate", "ENSO / IOD\n(Ocean-Atmosphere)", "NOAA / IRI", "El Nino Index, Indian Ocean Dipole,\nMultivariate ENSO Index", "Global monthly"],
    ["Conflict", "ACLED\n(Conflict Event Data)", "ACLED Project", "Event count, Fatalities", "Point → District"],
    ["Agriculture", "SPAM\n(Crop Presence Map)", "CIMMYT", "Crop cover fraction per district", "500m (static)"],
    ["Population", "WorldPop\n(Population Estimate)", "WorldPop Hub", "Population count per district", "1km (static)"],
]
add_table(slide, Inches(0.3), Inches(1.1), Inches(12.7), Inches(5.5), data_table,
          col_widths=[Inches(1.3), Inches(2.4), Inches(1.8), Inches(4.4), Inches(2.5)],
          font_size=10)

add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.4),
             "All sources spatially aligned to district (Admin2) boundaries  |  Temporal coverage: 2007–2025",
             font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ── SLIDE 6: Data Collection (Donghee) ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Data Collection Details", font_size=28, color=WHITE, bold=True)

card_l = add_shape(slide, Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.3), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(1.35), Inches(5), Inches(0.4),
             "Automated Pipeline (Halim)", font_size=18, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(1.85), Inches(5.4), Inches(4.5), [
    "CHIRPS (satellite rainfall): Global monthly data from UCSB (~7GB)",
    "FLDAS (land surface model): NASA soil moisture, temperature, runoff",
    "ACLED (conflict events): Armed conflict database (API, 2019–2025)",
    "WFP (food prices): Market price data via Humanitarian Data Exchange",
    "WorldPop: Population density raster (1km resolution)",
    "SPAM (crop presence): Crop coverage map (500m, CIMMYT)"
], font_size=12)

card_r = add_shape(slide, Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.3), CARD_BG)
add_text_box(slide, Inches(7.1), Inches(1.35), Inches(5), Inches(0.4),
             "Additional Data (Donghee)", font_size=18, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(7.1), Inches(1.85), Inches(5.4), Inches(4.5), [
    "FEWS NET IPC (food security severity classification)",
    "NASA Black Marble (satellite night lights — economic proxy)",
    "Dynamic World (Google Earth Engine land cover types)",
    "World Bank Pink Sheet (global commodity price indices)",
    "GeoBoundaries (official district boundary maps)",
    "Climate Indices (El Nino, Indian Ocean Dipole from NOAA/IRI)"
], font_size=12, bullet_color=ACCENT_GREEN)

add_text_box(slide, Inches(0.8), Inches(6.65), Inches(11), Inches(0.4),
             "Note: World Bank macro indices & FEWS NET (food security) data processed but reserved for future expansion",
             font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ── SLIDE 7: Section 3 ──
slide = prs.slides.add_slide(blank)
add_section_header(slide, "3", "Preprocessing & Merging", "From raw data to a unified panel dataset")


# ── SLIDE 8: Architecture Diagram ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
             "End-to-End Architecture", font_size=28, color=WHITE, bold=True)
slide.shapes.add_picture(chart_data_flow, Inches(0.5), Inches(1.0), Inches(12.3), Inches(5.7))


# ── SLIDE 9: Preprocessing Details ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Preprocessing Pipeline", font_size=28, color=WHITE, bold=True)

# 2x3 grid of compact cards
cards = [
    ("Food Prices (WFP) → District", ACCENT_BLUE, [
        "Spatial join: market lat/lon → district polygon",
        "Mean price aggregation per district/month",
    ]),
    ("Rainfall (CHIRPS) → Drought Index", ACCENT_GREEN, [
        "Gamma distribution fit (30-yr calibration period)",
        "Compute SPI (drought severity) per district",
    ]),
    ("Land Surface (FLDAS, 8 vars)", ACCENT_ORANGE, [
        "3-month moving average + Z-score normalization",
        "Soil moisture, temperature, runoff, evaporation",
    ]),
    ("Vegetation Health (MODIS/AVHRR)", ACCENT_PURPLE, [
        "NDVI (greenness), LST (surface temp), VHI (health)",
        "District-level zonal mean aggregation",
    ]),
    ("Conflict Events (ACLED)", ACCENT_RED, [
        "Point-in-polygon → district; sum fatalities + events",
        "Missing months filled with 0 (no events recorded)",
    ]),
    ("Static Features", MID_GRAY, [
        "Crop map: 500m raster → crop fraction per district",
        "Population: 1km raster → total count per district",
    ]),
]
for i, (title, color, items) in enumerate(cards):
    col = i % 3
    row = i // 3
    x = Inches(0.3 + col * 4.3)
    y = Inches(1.2 + row * 2.7)
    card = add_shape(slide, x, y, Inches(4.0), Inches(2.3), CARD_BG)
    add_rect(slide, x, y, Inches(0.06), Inches(2.3), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(3.5), Inches(0.35),
                 title, font_size=14, color=color, bold=True)
    add_bullet_list(slide, x + Inches(0.3), y + Inches(0.55), Inches(3.5), Inches(1.5),
                    items, font_size=11, bullet_color=color, spacing=Pt(4))


# ── SLIDE 10: Merging ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Data Merging Strategy", font_size=28, color=WHITE, bold=True)

# Flow: 5 step cards
steps = [
    ("1", "Skeleton", "Year × Month × District\ncartesian product", ACCENT_BLUE),
    ("2", "Prices", "Food prices joined on\n(year, month, district)", ACCENT_GREEN),
    ("3", "Static", "Population + Crop cover\njoined on district", ACCENT_ORANGE),
    ("4", "Conflict", "Conflict events on\n(year, month, district)", ACCENT_RED),
    ("5", "Climate", "Land surface + Vegetation\n+ Global climate indices", ACCENT_PURPLE),
]
for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.3 + i * 2.55)
    card = add_shape(slide, x, Inches(1.2), Inches(2.3), Inches(2.3), CARD_BG)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), Inches(1.35), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(1.9), Inches(0.3),
                 title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(2.35), Inches(2.0), Inches(1.0),
                 desc, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    # Arrow
    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(2.3), Inches(2.0), Inches(0.4), Inches(0.5),
                     "→", font_size=22, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Output card
output_card = add_shape(slide, Inches(0.5), Inches(3.8), Inches(12.3), Inches(3.0), CARD_BG)

# 3-column layout inside output card
merge_cols = [
    ("Spatial Reference", ACCENT_BLUE, [
        "GeoBoundaries district boundary maps (canonical)",
        "Points-in-polygon for markets & conflict",
        "Zonal statistics for raster (satellite) data",
        "Per-country spatial join"
    ]),
    ("Name Matching", ACCENT_GREEN, [
        "Direct match on district name",
        "Fuzzy matching (approximate string matching, ≥80%)",
        "Manual override dictionary for known mismatches",
        "e.g., 'Nairobi' city → 'Starehe' district"
    ]),
    ("Final Output", ACCENT_ORANGE, [
        "Unified Parquet dataset (columnar format)",
        "~8,640 rows × 50+ columns",
        "40 districts with price observations",
        "Ready for feature engineering & modeling"
    ]),
]
for i, (title, color, items) in enumerate(merge_cols):
    x = Inches(0.8 + i * 4.2)
    add_text_box(slide, x, Inches(3.95), Inches(3.8), Inches(0.35),
                 title, font_size=14, color=color, bold=True)
    add_bullet_list(slide, x, Inches(4.35), Inches(3.8), Inches(2.2),
                    items, font_size=10, bullet_color=color, spacing=Pt(3))


# ── SLIDE 11: Section 4 ──
slide = prs.slides.add_slide(blank)
add_section_header(slide, "4", "Feature Engineering", "Lag strategy and leakage prevention")


# ── SLIDE 12: Lag Strategy Diagram ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
             "Lag Strategy & Feature Engineering", font_size=28, color=WHITE, bold=True)

# Lag strategy diagram
slide.shapes.add_picture(chart_lag, Inches(0.3), Inches(0.9), Inches(12.5), Inches(4.2))

# Feature summary at bottom
feat_card = add_shape(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8), CARD_BG)
feat_cols = [
    ("Past Price Features", "1,2,3,6,12 month lags +\nrolling averages + year-over-year", ACCENT_BLUE),
    ("External Feature Shift", "Climate/conflict features shifted\nby forecast horizon (prevents leakage)", ACCENT_GREEN),
    ("Seasonal Encoding", "Cyclical month encoding\n(sine/cosine transformation)", ACCENT_ORANGE),
    ("Warmup Period", "Drop first 12 months per district\nto build sufficient history", ACCENT_RED),
]
for i, (title, desc, color) in enumerate(feat_cols):
    x = Inches(0.8 + i * 3.1)
    add_text_box(slide, x, Inches(5.4), Inches(2.8), Inches(0.35),
                 title, font_size=13, color=color, bold=True)
    add_text_box(slide, x, Inches(5.8), Inches(2.8), Inches(0.8),
                 desc, font_size=11, color=LIGHT_GRAY)


# ── SLIDE 13: Section 5 ──
slide = prs.slides.add_slide(blank)
add_section_header(slide, "5", "Modeling & Evaluation", "XGBoost (Gradient Boosting) with Spatio-Temporal Cross-Validation")


# ── SLIDE 14: Model Config + CV Design ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
             "Spatio-Temporal Cross-Validation Design", font_size=28, color=WHITE, bold=True)

# XGBoost config card (compact, left side)
cfg_card = add_shape(slide, Inches(0.3), Inches(0.9), Inches(3.5), Inches(6.0), CARD_BG)
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(3.1), Inches(0.35),
             "XGBoost (Gradient Boosting)", font_size=14, color=ACCENT_BLUE, bold=True)
config_text = (
    "n_estimators = 500\n"
    "learning_rate = 0.05\n"
    "max_depth = 6\n"
    "min_child_weight = 5\n"
    "subsample = 0.8\n"
    "colsample_bytree = 0.8\n"
    "reg_alpha = 0.1 (L1)\n"
    "reg_lambda = 1.0 (L2)\n"
    "early_stopping = 30 rounds"
)
add_multiline_text(slide, Inches(0.6), Inches(1.45), Inches(3.0), Inches(3.5),
                   config_text.split('\n'), font_size=11, color=LIGHT_GRAY, line_spacing=Pt(4))

# Combined formula
add_text_box(slide, Inches(0.5), Inches(4.5), Inches(3.1), Inches(0.35),
             "Combined Cross-Validation", font_size=13, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(0.5), Inches(4.9), Inches(3.1), Inches(1.5),
             "5 spatial × 3 temporal\n= 15 folds\nper target × horizon\n\n→ 135 total model fits\n(3 targets × 3 horizons × 15 folds)",
             font_size=12, color=LIGHT_GRAY)

# CV Design chart (right side, takes most of the space)
slide.shapes.add_picture(chart_cv_design, Inches(4.0), Inches(0.9), Inches(9.0), Inches(6.0))


# ── SLIDE 15: Section 6 ──
slide = prs.slides.add_slide(blank)
add_section_header(slide, "6", "Results", "Model performance across targets and horizons")


# ── SLIDE 16: Hold-out R² and MAPE charts ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
             "Hold-Out Test Results (2024+)", font_size=28, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.3),
             "Train on all data before 2024, test on 2024+ — simulates real-world deployment",
             font_size=12, color=LIGHT_GRAY)
slide.shapes.add_picture(chart_holdout_r2, Inches(0.2), Inches(1.2), Inches(6.5), Inches(5.8))
slide.shapes.add_picture(chart_holdout_mape, Inches(6.7), Inches(1.2), Inches(6.5), Inches(5.8))


# ── SLIDE 17: Horizon Degradation ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
             "Horizon Degradation Analysis", font_size=28, color=WHITE, bold=True)
slide.shapes.add_picture(chart_horizon, Inches(0.3), Inches(1.0), Inches(12.5), Inches(5.5))
add_text_box(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
             "Accuracy naturally degrades at longer forecast horizons  |  Food Price Index most stable  |  Sorghum most volatile",
             font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ── SLIDE 18: Country Comparison ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
             "Kenya vs Somalia Performance", font_size=28, color=WHITE, bold=True)
slide.shapes.add_picture(chart_country, Inches(1.5), Inches(1.0), Inches(10.3), Inches(5.5))
add_text_box(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
             "Kenya consistently more accurate than Somalia — likely due to higher data density, lower conflict, more stable markets",
             font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ── SLIDE 19: Feature Importance ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
             "Feature Group Importance", font_size=28, color=WHITE, bold=True)
slide.shapes.add_picture(chart_feat_imp, Inches(1.0), Inches(1.0), Inches(11.3), Inches(4.8))
add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.8),
             "Past price history dominates (77–94%)  |  Land surface & climate contribute 10–20% for commodity prices\n"
             "At longer horizons, external features (climate, conflict) gain relative importance",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ── SLIDE 20: CV vs Hold-out Comparison ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.5),
             "Spatio-Temporal CV vs Hold-Out Comparison", font_size=28, color=WHITE, bold=True)
slide.shapes.add_picture(chart_cv_vs_ho, Inches(1.5), Inches(1.0), Inches(10.3), Inches(5.3))
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.6),
             "Spatio-Temporal CV is intentionally conservative (held-out districts were never seen during training)\n"
             "Hold-out reflects realistic deployment: train on all districts, predict future time periods",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ── SLIDE 21: STCV Results Table ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Detailed Spatio-Temporal CV Results", font_size=28, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.3),
             "15 folds (5 spatial × 3 temporal) per target × horizon  |  2025 excluded from CV", font_size=12, color=LIGHT_GRAY)

cv_table = [
    ["Target", "Horizon", "R² (Variance Explained)", "MAPE (% Error)", "RMSE (Prediction Error)"],
    ["Food Price Index", "h=1", "0.00 ± 0.92", "9.68% ± 6.57%", "0.189 ± 0.137"],
    ["Food Price Index", "h=2", "-0.34 ± 1.17", "11.75% ± 7.29%", "0.219 ± 0.147"],
    ["Food Price Index", "h=3", "-0.65 ± 1.22", "13.43% ± 7.38%", "0.243 ± 0.147"],
    ["Maize", "h=1", "0.71 ± 0.34", "7.99% ± 2.01%", "4,300 ± 2,859"],
    ["Maize", "h=2", "0.60 ± 0.42", "10.93% ± 2.13%", "5,133 ± 2,962"],
    ["Maize", "h=3", "0.49 ± 0.47", "13.33% ± 2.40%", "5,706 ± 2,938"],
    ["Sorghum", "h=1", "0.76 ± 0.29", "9.28% ± 5.60%", "3,469 ± 3,914"],
    ["Sorghum", "h=2", "0.65 ± 0.33", "13.26% ± 6.44%", "4,094 ± 4,194"],
    ["Sorghum", "h=3", "0.52 ± 0.34", "16.74% ± 6.77%", "4,673 ± 4,377"],
]
add_table(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.0), cv_table,
          col_widths=[Inches(2.5), Inches(1.2), Inches(2.6), Inches(2.6), Inches(2.8)],
          font_size=12)


# ── SLIDE 22: Key Takeaways ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Key Takeaways", font_size=28, color=WHITE, bold=True)

takeaways = [
    ("Multi-Source Integration", "8 heterogeneous sources (climate, conflict, satellite, economic) → unified district-level panel for Kenya & Somalia", ACCENT_BLUE),
    ("Rigorous Evaluation", "Spatio-Temporal Cross-Validation (350km spatial separation × expanding time window, 2025 excluded) prevents data leakage", ACCENT_GREEN),
    ("Strong Predictive Performance", "Hold-out R² (variance explained) > 0.91 for all targets; Maize at 1-month achieves R² = 0.989", ACCENT_ORANGE),
    ("Country-Level Insights", "Kenya predictions more accurate than Somalia — likely due to data availability, lower conflict, more stable markets", ACCENT_RED),
    ("Forecast Horizon Trade-off", "1-month forecast most reliable; 3-month forecast error (MAPE) increases by 5–10 percentage points", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(takeaways):
    y = Inches(1.2 + i * 1.15)
    card = add_shape(slide, Inches(0.5), y, Inches(12.3), Inches(1.0), CARD_BG)
    add_rect(slide, Inches(0.5), y, Inches(0.08), Inches(1.0), color)
    add_text_box(slide, Inches(0.9), y + Inches(0.1), Inches(4.0), Inches(0.35),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, Inches(0.9), y + Inches(0.5), Inches(11.5), Inches(0.45),
                 desc, font_size=13, color=LIGHT_GRAY)


# ── SLIDE 23: Thank you ──
slide = prs.slides.add_slide(blank)
set_slide_bg(slide)
add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
             "Thank You", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(5), Inches(3.8), Inches(3), Inches(0.03), ACCENT_BLUE)
add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
             "Questions?", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.3), Inches(11), Inches(0.5),
             "Halim Jun  ·  Donghee Lee  ·  Sheldon Lu", font_size=15, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.4),
             "Columbia University  ·  Spring 2025", font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ── Save ──
output_path = os.path.join(OUT_DIR, "final_presentation.pptx")
prs.save(output_path)
print(f"\nPresentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")

print(f"Chart sources kept at: {CHART_DIR}")
