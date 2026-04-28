"""Insert charts from v2_final/_charts/ into v3 pptx at mapped slide positions."""

import os
from pptx import Presentation
from pptx.util import Inches

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
V3_IN = os.path.join(ROOT, "v3_final_v2", "final_presentationv2.pptx")
V3_OUT = os.path.join(ROOT, "v3_final_v2", "final_presentationv2_with_images.pptx")
CHARTS = os.path.join(ROOT, "v2_final", "_charts")

# v3 slide number (1-based, as shown in PowerPoint) -> list of (chart_filename, left_in, top_in, width_in, height_in)
INSERTIONS = {
    12: [("data_flow.png",          0.5, 1.0, 12.3, 5.7)],
    16: [("lag_strategy.png",       0.3, 0.9, 12.5, 4.2)],
    18: [("cv_design.png",          4.0, 0.9,  9.0, 6.0)],
    20: [("holdout_r2.png",         0.2, 1.2,  6.5, 5.8),
         ("holdout_mape.png",       6.7, 1.2,  6.5, 5.8)],
    21: [("horizon_degradation.png",0.3, 1.0, 12.5, 5.5)],
    22: [("country_comparison.png", 1.5, 1.0, 10.3, 5.5)],
    23: [("feature_importance.png", 1.0, 1.0, 11.3, 4.8)],
    24: [("cv_vs_holdout.png",      1.5, 1.0, 10.3, 5.3)],
}

prs = Presentation(V3_IN)
inserted = 0
for slide_num, items in INSERTIONS.items():
    slide = prs.slides[slide_num - 1]
    for fname, left, top, width, height in items:
        path = os.path.join(CHARTS, fname)
        if not os.path.exists(path):
            print(f"  [skip] missing: {path}")
            continue
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        print(f"  slide {slide_num}: inserted {fname}")
        inserted += 1

prs.save(V3_OUT)
print(f"\nInserted {inserted} images.")
print(f"Saved: {V3_OUT}")
